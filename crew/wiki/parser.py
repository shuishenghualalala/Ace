"""文档解析：把各种格式转换为原始 Markdown。"""

from __future__ import annotations

import re
import shutil
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Any

from crew.state.logging import get_logger

log = get_logger("wiki.parser")


class MissingDependencyError(RuntimeError):
    """缺少解析某种文档格式所需的可选依赖。"""

    def __init__(self, dependency: str, install_command: str, message: str) -> None:
        self.dependency = dependency
        self.install_command = install_command
        super().__init__(message)


class DocumentParseQualityError(RuntimeError):
    """文件虽然被解码，但结果不具备可入库的文本质量。"""


def _on_import_error(dependency: str, format_name: str, *, action: str = "解析") -> MissingDependencyError:
    """构建 MissingDependencyError，统一 install_command 和 message 格式。"""
    return MissingDependencyError(
        dependency=dependency,
        install_command='uv pip install -e ".[wiki]"',
        message=f"{action} {format_name} 需要 {dependency}，请安装: uv add {dependency}",
    )


def _decode_text_bytes(raw: bytes) -> str:
    if not raw:
        return ""
    if raw.startswith((b"\xff\xfe\x00\x00", b"\x00\x00\xfe\xff")):
        return raw.decode("utf-32")
    if raw.startswith((b"\xff\xfe", b"\xfe\xff")):
        return raw.decode("utf-16")
    if raw.startswith(b"\xef\xbb\xbf"):
        return raw.decode("utf-8-sig")
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        pass

    # Crew 主要面向中文环境；GB18030 是 GBK/GB2312 的超集。短中文文本对统计型
    # 探测器往往没有足够样本，会被误判为 Big5/韩文，因此先做严格 GB18030 解码。
    try:
        return raw.decode("gb18030")
    except UnicodeDecodeError:
        pass

    # charset-normalizer 是 Wiki 解析依赖；缺失时仍保留常见中文编码兜底。
    try:
        from charset_normalizer import from_bytes

        match = from_bytes(raw).best()
        if match is not None:
            return str(match)
    except ImportError:
        pass

    for encoding in ("big5",):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentParseQualityError("无法识别文本编码，请转换为 UTF-8 后重试")


def _looks_binary(raw: bytes) -> bool:
    if not raw or raw.startswith((
        b"\xff\xfe\x00\x00", b"\x00\x00\xfe\xff", b"\xff\xfe", b"\xfe\xff", b"\xef\xbb\xbf",
    )):
        return False
    binary_signatures = (
        b"\x89PNG\r\n\x1a\n",
        b"\xff\xd8\xff",
        b"GIF87a",
        b"GIF89a",
        b"PK\x03\x04",
        b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1",
        b"\x7fELF",
        b"MZ",
        b"\x1f\x8b",
    )
    if raw.startswith(binary_signatures):
        return True
    sample = raw[:8192]
    nul_ratio = sample.count(0) / len(sample)
    control_count = sum(1 for value in sample if value < 32 and value not in (9, 10, 13))
    return nul_ratio > 0.01 or control_count / len(sample) > 0.08


def validate_parsed_text(text: str, filename: str = "") -> str:
    """阻止空内容、乱码或二进制误解码结果进入后续 LLM 分析。"""
    stripped = text.strip()
    label = filename or "文件"
    if not stripped:
        raise DocumentParseQualityError(f"{label} 未提取到有效文本；可能是扫描件、空文档或受保护文件")
    length = max(1, len(text))
    replacement_ratio = text.count("\ufffd") / length
    control_ratio = sum(
        1 for char in text
        if ord(char) < 32 and char not in "\n\r\t"
    ) / length
    if replacement_ratio > 0.005:
        raise DocumentParseQualityError(
            f"{label} 解码后乱码比例过高（{replacement_ratio:.1%}），请检查文件格式或文本编码"
        )
    if "\x00" in text or control_ratio > 0.01:
        raise DocumentParseQualityError(f"{label} 包含大量二进制/控制字符，不能作为文本入库")
    return text


def _read_text(path: Path, *, errors: str = "strict") -> str:
    del errors  # 保留旧调用签名；解析质量门不再允许静默 replace。
    raw = path.read_bytes()
    if _looks_binary(raw):
        raise DocumentParseQualityError(f"不支持的二进制文件格式: {path.suffix.lower() or '未知'}")
    return _decode_text_bytes(raw)


def _markdown_table_from_rows(rows: list[list[str]]) -> str | None:
    """把二维行数据转成 Markdown 表格，过滤空行并统一列数。

    单元格内换行会被替换为空格，竖线会转义，避免破坏 Markdown 表格结构。
    """
    cleaned_rows: list[list[str]] = []
    max_cols = 0
    for row in rows:
        cells = [(c or "").strip().replace("\n", " ") for c in row]
        if not any(cells):
            continue
        cleaned_rows.append(cells)
        max_cols = max(max_cols, len(cells))

    if not cleaned_rows:
        return None

    for row in cleaned_rows:
        while len(row) < max_cols:
            row.append("")

    def _row_md(cells: list[str]) -> str:
        escaped = [c.replace("|", "\\|") for c in cells]
        return "| " + " | ".join(escaped) + " |"

    header = cleaned_rows[0]
    body = [r for r in cleaned_rows[1:] if any(cell.strip() for cell in r)]

    parts = [_row_md(header), "| " + " | ".join(["---"] * max_cols) + " |"]
    for row in body:
        parts.append(_row_md(row))
    return "\n".join(parts)


def _parse_pdf(path: Path) -> str:
    try:
        import fitz  # pymupdf
    except ImportError as exc:  # pragma: no cover
        raise _on_import_error("pymupdf", "PDF") from exc

    doc = fitz.open(str(path))
    parts: list[str] = []
    for i, page in enumerate(doc):
        page_parts: list[str] = []

        # 优先按块提取，保留段落/栏目结构，比 plain get_text 更不容易串行
        blocks = page.get_text("blocks")
        if blocks:
            block_texts: list[str] = []
            for b in blocks:
                # block tuple: (x0, y0, x1, y1, text, block_no, block_type)
                if len(b) >= 7 and b[6] == 0:
                    t = str(b[4]).strip()
                    if t:
                        block_texts.append(t)
            if block_texts:
                page_parts.append("\n\n".join(block_texts))
        else:
            text = page.get_text().strip()
            if text:
                page_parts.append(text)

        # 如果 PyMuPDF 支持，额外提取表格为 Markdown（避免纯文本块把表格拉乱）
        try:
            tables = page.find_tables()
            for table in tables:
                extracted = table.extract()
                if extracted:
                    rows = [[str(cell or "").strip() for cell in row] for row in extracted]
                    md = _markdown_table_from_rows(rows)
                    if md:
                        page_parts.append(md)
        except Exception:  # noqa: BLE001
            # 旧版 PyMuPDF 可能没有 find_tables，静默降级
            pass

        if page_parts:
            parts.append(f"## 第 {i + 1} 页\n\n" + "\n\n".join(page_parts))
    doc.close()
    return "\n\n".join(parts)


def _parse_docx(path: Path) -> str:
    try:
        import docx
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:  # pragma: no cover
        raise _on_import_error("python-docx", "DOCX") from exc

    document = docx.Document(str(path))
    parts: list[str] = []
    # python-docx 的 _Body 是 Paragraph/Table 的合法 parent；fake docx 测试可能没有 _body，
    # 此时回退到 document 本身以保持兼容。
    _parent = getattr(document, "_body", document)

    def _iter_story_elements(element):
        """按文档顺序产出段落/表格，递归进入内容控件/自定义 XML，跳过表格内部段落。

        Word 规范文档常把字段明细表放在 <w:sdt>（内容控件）或 <w:customXml> 中，
        仅遍历 body 直接子元素会丢失这些表格，因此需要递归展开。
        """
        for child in element:
            tag = child.tag
            if tag.endswith("}p"):
                yield ("p", Paragraph(child, _parent))
            elif tag.endswith("}tbl"):
                yield ("tbl", Table(child, _parent))
            elif tag.endswith("}sdt"):
                # 结构化文档标签（Content Control）— 进入 sdtContent 继续遍历
                sdt_content = child.find(qn("w:sdtContent"))
                if sdt_content is not None:
                    yield from _iter_story_elements(sdt_content)
            elif tag.endswith("}customXml"):
                # 自定义 XML 包装器 — 进入 customXmlContent 继续遍历
                cxn_content = child.find(qn("w:customXmlContent"))
                if cxn_content is not None:
                    yield from _iter_story_elements(cxn_content)

    for kind, obj in _iter_story_elements(document.element.body):
        if kind == "p":
            text = obj.text.strip()
            if text:
                parts.append(text)
        elif kind == "tbl":
            rows = [[cell.text for cell in row.cells] for row in obj.rows]
            md = _markdown_table_from_rows(rows)
            if md:
                parts.append(md)

    return "\n\n".join(parts)


# xlsx 读取上限：防止病态文件（整列刷格式产生的百万行/列声明）拖垮解析
# 以及后续 ingest 的 LLM 分块分析。
_XLSX_MAX_ROWS_PER_SHEET = 20_000
_XLSX_MAX_COLS = 200
_XLSX_MAX_CHARS = 500_000  # 整个工作簿输出字符上限


def _parse_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError as exc:  # pragma: no cover
        raise _on_import_error("openpyxl", "XLSX") from exc

    import shutil
    import zipfile
    from tempfile import TemporaryDirectory

    def _try_repair_xlsx(src: Path, dst: Path) -> bool:
        """在原有解析失败兜底链里插入：修复 Excel/WPS 生成的空 <fill/> 标签。"""
        tmp_dir = src.parent / f".{src.name}.repair.tmp"
        try:
            with zipfile.ZipFile(src, "r") as zin:
                zin.extractall(tmp_dir)

            styles_path = tmp_dir / "xl" / "styles.xml"
            if not styles_path.exists():
                return False

            text = styles_path.read_text(encoding="utf-8")
            repaired = re.sub(r"<fill\s*/>", '<fill><patternFill patternType="none"/></fill>', text)
            if repaired == text:
                return False

            styles_path.write_text(repaired, encoding="utf-8")
            log.warning("修复 xlsx 空 fill 标签: %s", src)

            with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
                for file_path in tmp_dir.rglob("*"):
                    if file_path.is_file():
                        zout.write(file_path, file_path.relative_to(tmp_dir))
            return True
        except Exception as exc:  # noqa: BLE001
            log.warning("修复 xlsx styles.xml 失败 %s: %s", src, exc)
            return False
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    with TemporaryDirectory() as tmp:
        # read_only 优先：流式读取，只占样式没有值的空单元格不会被物化。
        # 普通模式会把整列刷格式产生的海量空单元格全部建成对象，耗时和内存
        # 随单元格数线性爆炸（实测 50 万空样式单元格：4s vs 0.0s），是
        # 上传大 xlsx 卡死的主因。
        try:
            wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        except Exception:
            # 第一层兜底：修复已知的 Excel/WPS 样式 bug 后重试。
            repaired = Path(tmp) / "repaired.xlsx"
            if _try_repair_xlsx(path, repaired):
                try:
                    wb = openpyxl.load_workbook(str(repaired), data_only=True, read_only=True)
                except Exception:
                    log.warning("xlsx 修复后 read_only 仍失败，尝试普通模式读取: %s", path)
                    wb = openpyxl.load_workbook(str(repaired), data_only=True)
            else:
                # 没有可修复的已知问题，用普通模式做最后兜底。
                log.warning("xlsx read_only 解析失败，尝试普通模式读取: %s", path)
                wb = openpyxl.load_workbook(str(path), data_only=True)

        parts: list[str] = []
        remaining_chars = _XLSX_MAX_CHARS
        for sheet in wb.worksheets:
            parts.append(f"## 工作表: {sheet.title}")
            rows: list[list[str]] = []
            truncated = False
            used_cols = 0
            # 多取一行/一列用于判断是否还有数据被截断；
            # max_col 必须显式限制——dimension 声明成整行（16384 列）时
            # read_only 会把每行补齐到请求的宽度。
            for row in sheet.iter_rows(
                min_row=1,
                max_row=_XLSX_MAX_ROWS_PER_SHEET + 1,
                max_col=_XLSX_MAX_COLS + 1,
                values_only=True,
            ):
                if len(rows) >= _XLSX_MAX_ROWS_PER_SHEET:
                    truncated = True
                    break
                if len(row) > _XLSX_MAX_COLS:
                    if any(cell is not None for cell in row[_XLSX_MAX_COLS:]):
                        truncated = True
                    row = row[: _XLSX_MAX_COLS]
                cleaned = [str(cell) if cell is not None else "" for cell in row]
                row_chars = sum(len(c) for c in cleaned)
                if row_chars > remaining_chars:
                    truncated = True
                    break
                remaining_chars -= row_chars
                # 记录实际数据宽度，之后裁掉 read_only 补齐产生的尾部空列
                for i in range(len(cleaned) - 1, -1, -1):
                    if cleaned[i]:
                        used_cols = max(used_cols, i + 1)
                        break
                rows.append(cleaned)

            if not rows:
                parts.append("*(空工作表)*")
                continue

            rows = [r[:used_cols] for r in rows]
            md = _markdown_table_from_rows(rows)
            parts.append(md if md else "*(空工作表)*")
            if truncated:
                parts.append(f"*(内容过多，已截断：仅保留前 {len(rows)} 行、最多 {_XLSX_MAX_COLS} 列)*")

        result = "\n\n".join(parts)
        # read_only 模式持有 zip 文件句柄，显式关闭，避免 Windows 上
        # 临时目录清理失败及句柄泄漏。
        try:
            wb.close()
        except Exception:  # noqa: BLE001
            pass
        return result


_PPTX_IMAGE_PROMPT = (
    "请详细分析这张 PPT 中的图片，完成两件事："
    "1) 描述图片的内容、场景、对象、关系和在幻灯片中的作用；"
    "2) 转录图片中出现的所有可见文字（包括 UI 文字、标签、数字、标题等）。"
    "如果图片中没有文字，请明确说明'图中无文字'。"
)


def _shape_table_to_markdown(table: Any) -> str | None:
    """把 PPT 表格 shape 转为 Markdown 表格。"""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    return _markdown_table_from_rows(rows)


def _shape_chart_to_text(chart: Any) -> str | None:
    """把 PPT 图表 shape 转为文本摘要。"""
    try:
        chart_type = getattr(chart, "chart_type", None)
        chart_type_name = str(chart_type).replace("XL_CHART_TYPE.", "").replace("_", " ") if chart_type else "图表"

        series_summaries: list[str] = []
        for series in getattr(chart, "series", []):
            name = getattr(series, "name", "") or "系列"

            categories: list[str] = []
            try:
                cats = getattr(series, "categories", None)
                if cats:
                    categories = [
                        str(c.label) if hasattr(c, "label") else str(c) for c in cats
                    ]
            except Exception:  # noqa: BLE001
                pass

            values: list[str] = []
            try:
                values = [str(v) for v in getattr(series, "values", [])]
            except Exception:  # noqa: BLE001
                pass

            labels: list[str] = []
            try:
                if getattr(chart, "has_data_labels", False):
                    for point in getattr(series, "points", []):
                        try:
                            label = point.data_label.text_frame.text.strip()
                            if label:
                                labels.append(label)
                        except Exception:  # noqa: BLE001
                            pass
            except Exception:  # noqa: BLE001
                pass

            parts = [f'系列"{name}"']
            if categories:
                parts.append(f"类别{categories}")
            if values:
                parts.append(f"数值{values}")
            if labels:
                parts.append(f"数据标签{labels}")
            series_summaries.append("，".join(parts))

        if not series_summaries:
            return None
        return f"{chart_type_name}：" + "；".join(series_summaries)
    except Exception as exc:  # noqa: BLE001
        log.warning("PPT 图表解析失败: %s", exc)
        return None


def _extract_and_describe_image(
    shape: Any,
    tmp_dir: Path,
    *,
    failure_cache: dict[str, str] | None = None,
) -> str | None:
    """如果 shape 是图片，提取到临时目录并调用 VLM 描述+OCR。

    失败时返回用户可见的提示文本（而非 None），并把原因写入 failure_cache，
    避免同一次 PPT 解析中对每张图片重复发起注定失败的 VLM 请求。
    """
    cache = failure_cache if failure_cache is not None else {}
    if "reason" in cache:
        return f"[图片内容未解析：{cache['reason']}]"

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:  # noqa: BLE001
        return None

    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
        return None

    try:
        image = shape.image
        if not image or not image.blob:
            return None
    except Exception:  # noqa: BLE001
        return None

    ext = image.ext or "png"
    image_path = tmp_dir / f"pptx_image_{id(shape)}.{ext}"
    try:
        image_path.write_bytes(image.blob)
    except Exception:  # noqa: BLE001
        return None

    try:
        from crew.wiki.multimodal import MediaUnderstandingError, describe_image

        desc = describe_image(image_path, prompt=_PPTX_IMAGE_PROMPT)
        return desc.strip()
    except MediaUnderstandingError as exc:
        cache["reason"] = str(exc)
        log.warning("PPT 图片描述失败: %s", exc)
        return f"[图片内容未解析：{exc}]"
    except Exception as exc:  # noqa: BLE001
        cache["reason"] = str(exc)
        log.warning("PPT 图片描述失败: %s", exc)
        return f"[图片内容未解析：{exc}]"


def _parse_pptx(path: Path) -> str:
    try:
        import pptx
    except ImportError as exc:  # pragma: no cover
        raise _on_import_error("python-pptx", "PPTX") from exc

    from tempfile import TemporaryDirectory

    prs = pptx.Presentation(str(path))
    parts: list[str] = []
    with TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        image_failure_cache: dict[str, str] = {}
        for i, slide in enumerate(prs.slides, start=1):
            sections: list[str] = []

            # 布局名称
            try:
                layout_name = slide.slide_layout.name
                if layout_name:
                    sections.append(f"**布局**: {layout_name}")
            except Exception:  # noqa: BLE001
                pass

            text_parts: list[str] = []
            table_parts: list[str] = []
            chart_parts: list[str] = []
            image_parts: list[str] = []

            for shape in slide.shapes:
                # 表格
                if getattr(shape, "has_table", False):
                    table_md = _shape_table_to_markdown(shape.table)
                    if table_md:
                        table_parts.append(table_md)
                    continue

                # 图表
                if getattr(shape, "has_chart", False):
                    chart_text = _shape_chart_to_text(shape.chart)
                    if chart_text:
                        chart_parts.append(chart_text)
                    continue

                # 图片
                image_desc = _extract_and_describe_image(shape, tmp_dir, failure_cache=image_failure_cache)
                if image_desc:
                    image_parts.append(image_desc)

                # 文本框
                if getattr(shape, "has_text_frame", False):
                    texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                    if texts:
                        text_parts.append("\n".join(texts))

            if text_parts:
                sections.append("### 文本\n\n" + "\n\n".join(text_parts))
            if table_parts:
                sections.append("### 表格\n\n" + "\n\n".join(table_parts))
            if chart_parts:
                sections.append("### 图表\n\n" + "\n\n".join(chart_parts))
            if image_parts:
                sections.append("### 图片\n\n" + "\n\n".join(image_parts))

            # 演讲者备注
            try:
                if getattr(slide, "has_notes_slide", False):
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        sections.append(f"### 备注\n\n{notes_text}")
            except Exception:  # noqa: BLE001
                pass

            if sections:
                parts.append(f"## 第 {i} 页\n\n" + "\n\n".join(sections))
    return "\n\n".join(parts)


def _parse_html(path: Path) -> str:
    """将 HTML 文件转换为 Markdown。"""
    try:
        from markdownify import markdownify
    except ImportError as exc:  # pragma: no cover
        raise _on_import_error("markdownify", "HTML") from exc

    html_content = _decode_text_bytes(path.read_bytes())
    return markdownify(html_content, heading_style="ATX", bullets="-").strip()


def fetch_url_to_markdown(
    url: str,
    timeout: float = 15.0,
    allowed_targets: set[tuple[str, int, str]] | None = None,
) -> tuple[str, str]:
    """抓取 URL 并将 HTML 转为 Markdown。返回 (markdown_text, final_url)。"""
    from crew.security.outbound import fetch_public_http

    final_url, raw, content_type, charset = fetch_public_http(
        url,
        timeout=timeout,
        max_bytes=10_000_000,
        headers={"User-Agent": "Ace/1.0"},
        allowed_targets=allowed_targets,
    )

    if "html" in content_type:
        try:
            from markdownify import markdownify
        except ImportError as exc:  # pragma: no cover
            raise _on_import_error("markdownify", "HTML", action="转换") from exc
        html_str = raw.decode(charset, errors="replace")
        return markdownify(html_str, heading_style="ATX", bullets="-").strip(), final_url

    # 非 HTML 内容当纯文本
    return raw.decode(charset, errors="replace"), final_url


def _parse_legacy_office(path: Path, target_extension: str) -> str:
    """通过 LibreOffice 把旧版 Office 转成现代格式，再复用现有解析器。"""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise MissingDependencyError(
            dependency="LibreOffice",
            install_command="安装 LibreOffice，或先将文件另存为现代 Office 格式",
            message=(
                f"解析 {path.suffix.lower()} 需要 LibreOffice；"
                f"也可以先转换为 {target_extension} 后重新上传"
            ),
        )

    from tempfile import TemporaryDirectory

    with TemporaryDirectory(prefix=".ace-office-", dir=path.parent) as tmp:
        output_dir = Path(tmp)
        profile_dir = output_dir / "profile"
        profile_dir.mkdir()
        from crew.security.launch import execute_captured_sync

        try:
            completed = execute_captured_sync(
                (
                    soffice,
                    f"-env:UserInstallation={profile_dir.as_uri()}",
                    "--headless",
                    "--convert-to",
                    target_extension.lstrip("."),
                    "--outdir",
                    str(output_dir),
                    str(path),
                ),
                cwd=output_dir,
                timeout=120,
                tool_name="wiki_legacy_office",
            )
        except TimeoutError as exc:
            raise RuntimeError(f"{path.name} 通过 LibreOffice 转换超时") from exc
        converted = output_dir / f"{path.stem}{target_extension}"
        if completed.returncode != 0 or not converted.is_file():
            detail = (completed.stderr or completed.stdout or "").strip()
            raise RuntimeError(f"{path.name} 转换失败: {detail or 'LibreOffice 未生成输出文件'}")
        return parse_document_to_markdown(converted)


def parse_document_to_markdown(path: str | Path) -> str:
    """把文件解析为 Markdown 文本，并执行统一质量校验。"""
    p = Path(path)
    if not p.is_file():
        raise FileNotFoundError(f"文件不存在: {p}")

    ext = p.suffix.lower()
    if ext in (".html", ".htm"):
        result = _parse_html(p)
    else:
        text_extensions = {
            ".txt", ".md", ".markdown", ".json", ".py", ".js", ".ts", ".jsx", ".tsx",
            ".yaml", ".yml", ".css", ".xml", ".csv", ".log",
            ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp", ".cs", ".rb", ".php",
            ".swift", ".kt", ".scala", ".r", ".m", ".mm", ".sh", ".bash", ".zsh",
            ".ps1", ".sql", ".dockerfile", ".makefile", ".ini", ".cfg", ".conf",
            ".toml", ".llmwiki",  # 自定义 wiki 文本格式
        }
        if ext in text_extensions:
            result = _read_text(p)
        elif ext == ".pdf":
            result = _parse_pdf(p)
        elif ext == ".docx":
            result = _parse_docx(p)
        elif ext == ".xlsx":
            result = _parse_xlsx(p)
        elif ext == ".pptx":
            result = _parse_pptx(p)
        elif ext == ".doc":
            result = _parse_legacy_office(p, ".docx")
        elif ext == ".xls":
            result = _parse_legacy_office(p, ".xlsx")
        elif ext == ".ppt":
            result = _parse_legacy_office(p, ".pptx")
        else:
            # 保留无后缀/自定义后缀的真实文本支持；二进制会被 _read_text 明确拒绝。
            log.warning("未知文件类型 %s，尝试识别为文本", ext)
            result = _read_text(p)
    return validate_parsed_text(result, p.name)


def _detect_content_extension(content: bytes) -> str:
    """从可靠文件头识别可解析的文档格式；无法判断时返回空字符串。"""
    if content.startswith(b"%PDF-"):
        return ".pdf"
    if content.startswith(b"\x89PNG\r\n\x1a\n"):
        return ".png"
    if content.startswith(b"\xff\xd8\xff"):
        return ".jpg"
    if content.startswith((b"GIF87a", b"GIF89a")):
        return ".gif"
    if content.startswith(b"BM"):
        return ".bmp"
    if content.startswith(b"RIFF") and content[8:12] == b"WEBP":
        return ".webp"
    if content.startswith(b"RIFF") and content[8:12] == b"AVI ":
        return ".avi"
    if len(content) >= 12 and content[4:8] == b"ftyp":
        return ".mp4"
    if content.startswith(b"PK\x03\x04"):
        try:
            with zipfile.ZipFile(BytesIO(content)) as archive:
                names = set(archive.namelist())
        except (OSError, zipfile.BadZipFile):
            return ""
        if any(name.startswith("word/") for name in names):
            return ".docx"
        if any(name.startswith("xl/") for name in names):
            return ".xlsx"
        if any(name.startswith("ppt/") for name in names):
            return ".pptx"
    return ""


def _document_extension(filename: str | Path, content: bytes | None = None) -> str:
    """确定文档扩展名：可靠文件头优先，其次才使用文件名后缀。"""
    detected = _detect_content_extension(content) if content is not None else ""
    return detected or Path(filename).suffix.lower()


def guess_mime_type(path: str | Path, content: bytes | None = None) -> str:
    """根据文件头和后缀猜测 MIME 类型。"""
    ext = _document_extension(path, content)
    mapping: dict[str, str] = {
        ".html": "text/html",
        ".htm": "text/html",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt": "application/vnd.ms-powerpoint",
        ".json": "application/json",
        ".yaml": "application/yaml",
        ".yml": "application/yaml",
        # 图片
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".webp": "image/webp",
        ".bmp": "image/bmp",
        ".gif": "image/gif",
        # 视频
        ".mp4": "video/mp4",
        ".mov": "video/quicktime",
        ".webm": "video/webm",
        ".avi": "video/x-msvideo",
        ".mkv": "video/x-matroska",
        ".flv": "video/x-flv",
        ".m4v": "video/x-m4v",
    }
    return mapping.get(ext, "application/octet-stream")


def parse_document_from_bytes(content: bytes, filename: str) -> str:
    """从内存字节解析文档（临时写文件后解析）。"""
    from tempfile import TemporaryDirectory
    from crew.security.launch import current_process_launch

    ext = _document_extension(filename, content) or ".txt"
    launch = current_process_launch.get()
    workspace_root = launch.security_context.workspace_root if launch and launch.security_context else None
    temporary_parent = workspace_root if workspace_root and workspace_root.is_dir() else None
    with TemporaryDirectory(prefix=".ace-wiki-", dir=temporary_parent) as tmp:
        tmp_path = Path(tmp) / f"upload{ext}"
        tmp_path.write_bytes(content)
        return parse_document_to_markdown(tmp_path)
